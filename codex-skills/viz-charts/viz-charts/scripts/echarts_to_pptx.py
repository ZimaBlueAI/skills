#!/usr/bin/env python
"""viz-charts → native PPTX chart bridge.

Input: chart spec JSON (subset of ECharts option, plus title/canvas).
Output: PPTX with a single slide containing a native data-bound chart.

The chart is a real `<c:chart>` object inside the PPTX — opening it in
PowerPoint shows the underlying data in a spreadsheet panel and lets the
user edit series, recolor, change chart type, all natively.

Run with the ppt-master venv:
  ~/.claude/skills/ppt-master/.venv/Scripts/python.exe \
    ~/.claude/skills/viz-charts/scripts/echarts_to_pptx.py spec.json -o chart.pptx

Spec schema:
{
  "title":      "Quarterly Revenue",
  "subtitle":   "Across all regions, FY 2026",
  "type":       "bar" | "column" | "line" | "pie" | "doughnut" | "area" | "scatter" | "radar",
  "stacked":    true | false,                # optional, for bar/column/area/line
  "categories": ["Q1", "Q2", "Q3", "Q4"],    # x-axis labels (omitted for pie/doughnut)
  "series": [
    { "name": "Americas", "data": [12, 18, 22, 28], "color": "#00D9FF" },
    { "name": "EMEA",     "data": [ 8, 11, 14, 18], "color": "#6366F1" },
    { "name": "APAC",     "data": [ 5,  7,  9, 13], "color": "#FFB800" }
  ],
  "y_axis":     "USD millions",              # optional
  "x_axis":     "Quarter",                   # optional
  "canvas":     "ppt169" | "ppt43",          # optional, default ppt169
  "theme":      "deep-space" | "terminal" | "deck-light"
}
"""

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor

# canvas formats (EMU = 914400 per inch)
CANVAS = {
    "ppt169": (13333333, 7500000),  # 14.81in x 8.33in -> 1280x720 @ 96dpi mapping
    "ppt43": (9144000, 6858000),    # 10in x 7.5in
}

THEMES = {
    "deep-space": {
        "bg": RGBColor(0x03, 0x07, 0x11),
        "text": RGBColor(0xE5, 0xE7, 0xEB),
        "muted": RGBColor(0x94, 0xA3, 0xB8),
        "accent": RGBColor(0x00, 0xD9, 0xFF),
        "series_palette": [
            RGBColor(0x00, 0xD9, 0xFF), RGBColor(0x63, 0x66, 0xF1),
            RGBColor(0xFF, 0xB8, 0x00), RGBColor(0xEC, 0x40, 0x99),
            RGBColor(0x10, 0xB9, 0x81), RGBColor(0xF9, 0x73, 0x16),
        ],
    },
    "terminal": {
        "bg": RGBColor(0x0A, 0x0A, 0x0A),
        "text": RGBColor(0xF5, 0xF5, 0xF5),
        "muted": RGBColor(0x9C, 0xA3, 0xAF),
        "accent": RGBColor(0xD4, 0xFF, 0x00),
        "series_palette": [
            RGBColor(0xD4, 0xFF, 0x00), RGBColor(0xFF, 0xFF, 0xFF),
            RGBColor(0xFF, 0xB8, 0x00), RGBColor(0x9C, 0xA3, 0xAF),
            RGBColor(0xEF, 0x44, 0x44), RGBColor(0x10, 0xB9, 0x81),
        ],
    },
    "deck-light": {
        "bg": RGBColor(0xFA, 0xFB, 0xFC),
        "text": RGBColor(0x0F, 0x17, 0x2A),
        "muted": RGBColor(0x47, 0x55, 0x69),
        "accent": RGBColor(0x63, 0x66, 0xF1),
        "series_palette": [
            RGBColor(0x63, 0x66, 0xF1), RGBColor(0x06, 0xB6, 0xD4),
            RGBColor(0xF5, 0x9E, 0x0B), RGBColor(0x10, 0xB9, 0x81),
            RGBColor(0xEF, 0x44, 0x44), RGBColor(0x8B, 0x5C, 0xF6),
        ],
    },
}

# ECharts type → python-pptx chart type
CHART_TYPE = {
    ("bar", False): XL_CHART_TYPE.BAR_CLUSTERED,
    ("bar", True):  XL_CHART_TYPE.BAR_STACKED,
    ("column", False): XL_CHART_TYPE.COLUMN_CLUSTERED,
    ("column", True):  XL_CHART_TYPE.COLUMN_STACKED,
    ("line", False): XL_CHART_TYPE.LINE,
    ("line", True):  XL_CHART_TYPE.LINE_STACKED,
    ("area", False): XL_CHART_TYPE.AREA,
    ("area", True):  XL_CHART_TYPE.AREA_STACKED,
    ("pie", False): XL_CHART_TYPE.PIE,
    ("doughnut", False): XL_CHART_TYPE.DOUGHNUT,
    ("scatter", False): XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,
    ("radar", False): XL_CHART_TYPE.RADAR,
}


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    if len(h) != 6:
        return None
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def normalize(spec: dict) -> dict:
    """Accept either our flat spec OR a raw ECharts option object."""
    if "series" in spec and ("categories" in spec or spec.get("type") in ("pie", "doughnut", "scatter")):
        return spec
    # Try ECharts option object
    if "series" in spec and "xAxis" in spec:
        x = spec["xAxis"]
        if isinstance(x, list):
            x = x[0]
        categories = x.get("data", [])
        series = spec["series"]
        s_type = series[0].get("type", "line") if series else "line"
        if s_type in ("bar",):
            s_type = "column"  # ECharts "bar" with vertical orientation == PPT column
        out = {
            "title": (spec.get("title") or {}).get("text") if isinstance(spec.get("title"), dict) else spec.get("title"),
            "type": s_type,
            "stacked": any(s.get("stack") for s in series),
            "categories": categories,
            "series": [
                {"name": s.get("name", f"Series {i+1}"), "data": s.get("data", [])}
                for i, s in enumerate(series)
            ],
        }
        return out
    return spec


def add_native_chart(slide, spec: dict, theme: dict):
    chart_type_key = (spec.get("type", "column").lower(), bool(spec.get("stacked", False)))
    if chart_type_key not in CHART_TYPE:
        # Fall back to non-stacked variant
        chart_type_key = (chart_type_key[0], False)
    chart_type = CHART_TYPE.get(chart_type_key, XL_CHART_TYPE.COLUMN_CLUSTERED)

    if spec.get("type") == "scatter":
        chart_data = XyChartData()
        for i, s in enumerate(spec.get("series", [])):
            series_obj = chart_data.add_series(s.get("name", f"Series {i+1}"))
            for point in s.get("data", []):
                x, y = (point[0], point[1]) if isinstance(point, (list, tuple)) else (i, point)
                series_obj.add_data_point(x, y)
    else:
        chart_data = CategoryChartData()
        chart_data.categories = spec.get("categories", [])
        for s in spec.get("series", []):
            chart_data.add_series(s.get("name", "Series"), s.get("data", []))

    # Place chart on slide
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(12.3) if spec.get("canvas", "ppt169") == "ppt169" else Inches(9.0)
    height = Inches(5.8)
    graphic_frame = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
    chart = graphic_frame.chart

    # Style: theme palette + legend
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # Apply series colors
    palette = theme["series_palette"]
    for i, plot_series in enumerate(chart.series):
        spec_series = spec.get("series", [])[i] if i < len(spec.get("series", [])) else {}
        color_hex = spec_series.get("color")
        color = hex_to_rgb(color_hex) if color_hex else palette[i % len(palette)]
        if color is None:
            continue
        fill = plot_series.format.fill
        fill.solid()
        fill.fore_color.rgb = color
        line = plot_series.format.line
        line.color.rgb = color

    return chart


def add_title_block(slide, spec: dict, theme: dict, canvas_w_emu: int):
    from pptx.util import Emu
    # Title (top-left)
    if spec.get("title"):
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Emu(canvas_w_emu - Inches(1).emu), Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = spec["title"]
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = theme["text"]
    if spec.get("subtitle"):
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Emu(canvas_w_emu - Inches(1).emu), Inches(0.3))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = spec["subtitle"]
        run.font.size = Pt(14)
        run.font.color.rgb = theme["muted"]


def set_slide_background(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def build(spec: dict, out_path: Path):
    canvas = spec.get("canvas", "ppt169")
    if canvas not in CANVAS:
        canvas = "ppt169"
    w, h = CANVAS[canvas]
    theme = THEMES.get(spec.get("theme", "deep-space"), THEMES["deep-space"])

    prs = Presentation()
    prs.slide_width = Emu(w)
    prs.slide_height = Emu(h)
    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)

    set_slide_background(slide, theme["bg"])
    add_title_block(slide, spec, theme, w)
    add_native_chart(slide, spec, theme)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"[viz-charts] OK native PPTX chart: {out_path}")


def main():
    ap = argparse.ArgumentParser(description="Convert chart spec (or ECharts option) → editable PPTX with native chart")
    ap.add_argument("spec", help="JSON spec file")
    ap.add_argument("-o", "--out", required=True, help="Output .pptx path")
    args = ap.parse_args()

    spec_path = Path(args.spec)
    if not spec_path.exists():
        print(f"ERROR: spec file not found: {spec_path}", file=sys.stderr)
        sys.exit(2)

    spec = normalize(json.loads(spec_path.read_text(encoding="utf-8")))
    build(spec, Path(args.out))


if __name__ == "__main__":
    main()
